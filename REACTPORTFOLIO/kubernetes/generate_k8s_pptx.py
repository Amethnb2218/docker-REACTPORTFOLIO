from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2 = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK


def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)
    return slide


def section_slide(num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3), Inches(0.12), Inches(1.5))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(3), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(10), Inches(1.2))
    p2 = tx2.text_frame.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    return slide


def content_slide(title, bullets, max_items=7, color=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = color
    tx2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.5), Inches(5.5))
    tf = tx2.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets[:max_items]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = b
        para.font.size = Pt(17)
        para.font.color.rgb = WHITE if not b.startswith("   ") else GRAY
        para.space_after = Pt(10)
    return slide


def code_slide(title, code):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x13, 0x13, 0x25)
    sh.line.color.rgb = RGBColor(0x2D, 0x2D, 0x44)
    sh.line.width = Pt(1)
    tx2 = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.4))
    tf = tx2.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    p2.text = code
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    p2.font.name = "Consolas"
    return slide


def table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.8)
    top = Inches(1.3)
    width = Inches(11.5)
    row_h = min(Inches(0.7), Inches(5.5) / n_rows)
    height = row_h * n_rows
    table = slide.shapes.add_table(n_rows, cols, left, top, width, height).table
    for i in range(cols):
        table.columns[i].width = int(width / cols)
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x40, 0x6E)
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(13)
            para.font.bold = True
            para.font.color.rgb = WHITE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B) if r % 2 == 0 else BG_DARK
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.color.rgb = WHITE
    return slide


# ============ SLIDES ============

title_slide("Kubernetes", "Orchestration de conteneurs\nProjet fil rouge - Portfolio React Express\nGroupe 5")

content_slide("Sommaire", [
    "5.1  Problematique",
    "5.2  Presentation",
    "5.3  Concepts",
    "5.4  Architecture d'un cluster",
    "5.5  Communication (interne / externe)",
    "5.6  Installation",
    "5.7  Commandes de base",
    "5.8  Demo",
    "5.9  References",
])

# 5.1
section_slide("5.1", "Problematique")

content_slide("Limites de Docker Compose", [
    "Pas de haute disponibilite",
    "Pas de scaling automatique",
    "Deploiement mono-serveur uniquement",
    "Mises a jour avec temps d'arret",
    "Secrets en clair dans les fichiers",
    "Pas de self-healing",
])

table_slide("Ce que Kubernetes resout",
    ["Probleme", "Docker Compose", "Kubernetes"],
    [
        ["Crash conteneur", "Intervention manuelle", "Auto-restart"],
        ["Scalabilite", "Manuelle", "HPA automatique"],
        ["Haute dispo", "Non", "ReplicaSets multi-noeuds"],
        ["Mise a jour", "Downtime", "Rolling update"],
        ["Secrets", "En clair", "Objets chiffres + RBAC"],
    ])

# 5.2
section_slide("5.2", "Presentation")

content_slide("Qu'est-ce que Kubernetes ?", [
    "Plateforme open-source d'orchestration de conteneurs",
    "Cree par Google, open-source depuis 2014",
    "Maintenu par la CNCF",
    "Automatise deploiement, scaling et gestion",
    "Standard cloud-native de l'industrie",
])

content_slide("Fonctionnalites cles", [
    "Self-healing : redemarrage et remplacement auto",
    "Auto-scaling : HPA, VPA, Cluster Autoscaler",
    "Service Discovery + Load Balancing",
    "Rolling Updates + Rollback instantane",
    "Configuration declarative (YAML)",
    "Extensible via CRDs et Operators",
])

# 5.3
section_slide("5.3", "Concepts")

content_slide("Node", [
    "Machine (physique ou VM) dans le cluster",
    "",
    "Master Node (Control Plane) :",
    "   API Server, Scheduler, Controller Manager, etcd",
    "",
    "Worker Node :",
    "   kubelet, kube-proxy, Container Runtime, Pods",
])

content_slide("Pod", [
    "Plus petite unite deployable",
    "Contient 1+ conteneurs (meme reseau/stockage)",
    "Adresse IP unique dans le cluster",
    "Ephemere : peut etre detruit/recree",
])

content_slide("ReplicaSet & Deployment", [
    "ReplicaSet : maintient N pods actifs",
    "Deployment : gere ReplicaSet + updates",
    "   - Rolling updates (zero downtime)",
    "   - Historique des revisions",
    "   - Rollback en 1 commande",
    "",
    "Notre projet : 2 replicas frontend + 2 backend",
])

content_slide("Service", [
    "IP stable pour un ensemble de pods",
    "",
    "ClusterIP : interne au cluster (defaut)",
    "NodePort : port sur chaque noeud",
    "LoadBalancer : IP publique cloud",
    "",
    "DNS : backend-service.portfolio.svc.cluster.local",
])

content_slide("Ingress", [
    "Routage HTTP/HTTPS externe",
    "Routage par path ou hostname",
    "Terminaison TLS",
    "Necessite un Ingress Controller (nginx)",
    "",
    "portfolio.local/     -> frontend",
    "portfolio.local/api  -> backend",
])

content_slide("ConfigMap, Secret, Volume, StatefulSet", [
    "ConfigMap : config non-sensible (PORT, NODE_ENV)",
    "Secret : donnees sensibles (MONGO_URI, passwords)",
    "",
    "PersistentVolume : stockage durable",
    "PersistentVolumeClaim : demande de stockage",
    "",
    "StatefulSet : apps stateful (MongoDB)",
    "   - Identite stable, stockage persistant par pod",
])

# 5.4
section_slide("5.4", "Architecture")

code_slide("Control Plane",
"""CONTROL PLANE (Master)
======================

┌─────────────┐  ┌────────────┐  ┌───────────────────┐  ┌────────┐
│  API Server │  │  Scheduler │  │ Controller Manager │  │  etcd  │
│             │  │            │  │                   │  │        │
│ Expose API  │  │ Place pods │  │ Boucle controle   │  │ State  │
│ Valide obj  │  │ sur noeuds │  │ Maintient desired │  │ store  │
└─────────────┘  └────────────┘  └───────────────────┘  └────────┘""")

code_slide("Architecture du projet",
"""           Namespace: portfolio
           ======================

Internet --> Ingress (nginx)
                 |
         +-------+--------+
         |                |
         v                v
  frontend-service   backend-service
      :80                :5000
         |                |
    +----+----+      +----+----+
    |         |      |         |
   Pod F1   Pod F2  Pod B1   Pod B2 --> mongo-service:27017
                                              |
  Deployment(2)    Deployment(2)         StatefulSet
                                          + PVC 2Gi""")

# 5.5
section_slide("5.5", "Communication")

content_slide("Communication interne", [
    "DNS CoreDNS : <service>.<namespace>.svc.cluster.local",
    "Pod-to-Pod : IP unique, CNI plugin (Calico, Flannel)",
    "Pod-to-Service : ClusterIP + kube-proxy",
    "",
    "Notre projet :",
    "   Frontend -> backend-service:5000",
    "   Backend -> mongo-service:27017",
])

content_slide("Communication externe", [
    "NodePort : port statique sur chaque noeud",
    "LoadBalancer : cloud provider (AWS ELB, GCP LB)",
    "Ingress : routage HTTP L7 (notre choix)",
    "",
    "NetworkPolicy : securite reseau",
    "   MongoDB accessible UNIQUEMENT par le backend",
])

# 5.6
section_slide("5.6", "Installation")

content_slide("Options locales", [
    "Docker Desktop : Settings > Kubernetes > Enable",
    "Minikube : minikube start --driver=docker",
    "Kind : kind create cluster",
    "",
    "kubectl : choco install kubernetes-cli",
    "Ingress : minikube addons enable ingress",
])

code_slide("Deployer le projet",
"""# Build images
docker build -t portfolio-frontend ./reactportfolio
docker build -t portfolio-backend "./EXPRESSJS PORTFOLIO"

# Charger dans Minikube
minikube image load portfolio-frontend:latest
minikube image load portfolio-backend:latest

# Deployer
cd kubernetes
./deploy.sh up

# Ajouter dans /etc/hosts
127.0.0.1 portfolio.local

# Acceder : http://portfolio.local""")

# 5.7
section_slide("5.7", "Commandes de base")

code_slide("Commandes essentielles",
"""kubectl get nodes                           # Noeuds du cluster
kubectl get pods -n portfolio               # Pods du namespace
kubectl get svc -n portfolio                # Services
kubectl get all -n portfolio                # Tout

kubectl describe pod <name> -n portfolio    # Details
kubectl logs -f <pod> -n portfolio          # Logs temps reel
kubectl exec -it <pod> -n portfolio -- sh   # Shell dans pod

kubectl scale deployment backend --replicas=4 -n portfolio
kubectl rollout status deployment/backend -n portfolio
kubectl rollout undo deployment/backend -n portfolio

kubectl port-forward svc/frontend-service 3000:80 -n portfolio""")

# 5.8
section_slide("5.8", "Demo")

content_slide("Ce qui est deploye", [
    "Namespace : portfolio",
    "Frontend : 2 pods (Deployment + Service + HPA)",
    "Backend : 2 pods (Deployment + Service + HPA)",
    "MongoDB : StatefulSet + PVC 2Gi",
    "Ingress nginx : portfolio.local",
    "NetworkPolicy : acces mongo restreint",
])

content_slide("Demo 1 : Self-healing", [
    "kubectl delete pod -l app=backend -n portfolio",
    "",
    "-> Le pod est recree automatiquement",
    "-> Zero intervention, l'app reste disponible",
    "-> Le ReplicaSet maintient toujours 2 pods",
])

content_slide("Demo 2 : Scaling + Rolling Update", [
    "# Scaler",
    "kubectl scale deployment backend --replicas=4",
    "",
    "# Update",
    "kubectl set image deployment/backend backend=v2",
    "-> Mise a jour progressive, zero downtime",
    "",
    "# Rollback",
    "kubectl rollout undo deployment/backend",
])

table_slide("Docker Compose -> Kubernetes",
    ["Docker Compose", "Kubernetes", "Fichier"],
    [
        ["services.frontend", "Deployment + Service", "frontend-deployment.yaml"],
        ["services.mongo", "StatefulSet + Headless Svc", "mongo-statefulset.yaml"],
        ["volumes", "PV + PVC", "mongo-pv.yaml"],
        ["environment", "ConfigMap + Secret", "configmap/secret.yaml"],
        ["ports", "Ingress", "ingress.yaml"],
        ["(aucun)", "HPA auto-scaling", "hpa.yaml"],
    ])

# 5.9
section_slide("5.9", "References")

content_slide("Sources", [
    "kubernetes.io/docs/ - Documentation officielle",
    "kubernetes.io/docs/reference/kubectl/cheatsheet/",
    "gateway-api.sigs.k8s.io/ - Gateway API",
    "minikube.sigs.k8s.io - Minikube",
    "helm.sh - Helm (package manager)",
    "killercoda.com/kubernetes - Tutoriels interactifs",
])

title_slide("Merci !", "Questions ?")

# Save
out = r"C:\Users\khadjia\Desktop\docker\REACTPORTFOLIO\kubernetes\Kubernetes_Presentation.pptx"
prs.save(out)
print(f"OK: {out} ({len(prs.slides)} slides)")
