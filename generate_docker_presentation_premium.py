# -*- coding: utf-8 -*-
from pathlib import Path
import math
import random

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from generate_docker_presentation import SLIDES


OUT = Path("Docker_Cours_Premium_90_slides.pptx")
ASSET_DIR = Path("premium_pptx_assets")
W, H = 1920, 1080


def rgb(value):
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def hex_tuple(value):
    value = value.strip("#")
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def mix(a, b, t):
    return tuple(int(a[i] * (1 - t) + b[i] * t) for i in range(3))


COLORS = {
    "ink": rgb("111827"),
    "ink2": rgb("243B53"),
    "muted": rgb("64748B"),
    "soft": rgb("EEF4FA"),
    "line": rgb("D7E1EA"),
    "white": rgb("FFFFFF"),
    "night": rgb("071827"),
    "code": rgb("0B1220"),
    "blue": rgb("2563EB"),
    "cyan": rgb("0891B2"),
    "teal": rgb("0F766E"),
    "green": rgb("15803D"),
    "amber": rgb("D97706"),
    "orange": rgb("EA580C"),
    "red": rgb("DC2626"),
    "purple": rgb("7C3AED"),
    "slate": rgb("475569"),
}


ACCENTS = {
    "2.1": ("#EA580C", "#F59E0B"),
    "2.2": ("#0F766E", "#22C55E"),
    "2.3": ("#2563EB", "#06B6D4"),
    "2.4": ("#0F172A", "#2563EB"),
    "2.5": ("#0891B2", "#38BDF8"),
    "2.6": ("#7C3AED", "#06B6D4"),
    "2.7": ("#15803D", "#84CC16"),
    "2.8": ("#D97706", "#FACC15"),
    "2.9": ("#475569", "#0EA5E9"),
    "2. Docker": ("#2563EB", "#22D3EE"),
}


SECTION_LABELS = {
    "2.1": "Pourquoi Docker ?",
    "2.2": "Conteneurisation",
    "2.3": "Installation",
    "2.4": "Commandes de base",
    "2.5": "Réseaux",
    "2.6": "Images Docker",
    "2.7": "Docker Hub",
    "2.8": "Stockage",
    "2.9": "Docker Compose",
    "2. Docker": "Vue d’ensemble",
}


def section_key(section):
    for key in ["2.1", "2.2", "2.3", "2.4", "2.5", "2.6", "2.7", "2.8", "2.9"]:
        if section.startswith(key):
            return key
    return "2. Docker"


def accent_pair(section):
    return ACCENTS[section_key(section)]


def accent_rgb(section):
    return rgb(accent_pair(section)[0])


def add_noise(img, alpha=10):
    rnd = random.Random(42)
    draw = ImageDraw.Draw(img, "RGBA")
    for _ in range(1800):
        x = rnd.randrange(W)
        y = rnd.randrange(H)
        v = rnd.randint(255 - alpha, 255)
        draw.point((x, y), fill=(v, v, v, 18))


def draw_grid(draw, color, step=64, alpha=22):
    rgba = color + (alpha,)
    for x in range(0, W, step):
        draw.line([(x, 0), (x, H)], fill=rgba, width=1)
    for y in range(0, H, step):
        draw.line([(0, y), (W, y)], fill=rgba, width=1)


def draw_glow(base, center, radius, color, opacity=150):
    layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(layer, "RGBA")
    x, y = center
    d.ellipse((x - radius, y - radius, x + radius, y + radius), fill=color + (opacity,))
    layer = layer.filter(ImageFilter.GaussianBlur(radius // 2))
    base.alpha_composite(layer)


def make_gradient(c1, c2):
    strip = Image.new("RGB", (1, H), c1)
    px = strip.load()
    for y in range(H):
        t = y / (H - 1)
        px[0, y] = mix(c1, c2, t)
    img = strip.resize((W, H))
    side = Image.new("RGBA", (W, H), (255, 255, 255, 0))
    d = ImageDraw.Draw(side, "RGBA")
    for x in range(0, W, 24):
        alpha = int(32 * x / W)
        d.rectangle((x, 0, x + 24, H), fill=(255, 255, 255, alpha))
    img = img.convert("RGBA")
    img.alpha_composite(side)
    return img.convert("RGBA")


def section_slug(section):
    return section_key(section).replace(".", "_").replace(" ", "_")


def create_light_background(section):
    a1, a2 = (hex_tuple(v) for v in accent_pair(section))
    img = make_gradient(hex_tuple("#F9FBFD"), hex_tuple("#ECF3F9"))
    draw = ImageDraw.Draw(img, "RGBA")
    draw_grid(draw, hex_tuple("#8AA4BB"), step=72, alpha=18)
    draw_glow(img, (1660, 245), 440, a1, 75)
    draw_glow(img, (1780, 890), 380, a2, 55)

    for i in range(9):
        y = 140 + i * 84
        x0 = 1310 + int(math.sin(i * 0.8) * 36)
        draw.line([(x0, y), (1850, y + 130)], fill=a1 + (30,), width=2)

    for i in range(5):
        draw.rounded_rectangle(
            (1450 + i * 54, 150 + i * 34, 1760 + i * 54, 204 + i * 34),
            radius=16,
            outline=a1 + (42,),
            width=3,
        )

    draw.rounded_rectangle((84, 88, 1836, 992), radius=44, outline=(255, 255, 255, 128), width=2)
    add_noise(img, alpha=3)
    return img.convert("RGB")


def create_dark_background(section):
    a1, a2 = (hex_tuple(v) for v in accent_pair(section))
    img = make_gradient(hex_tuple("#071421"), hex_tuple("#0B2236"))
    draw = ImageDraw.Draw(img, "RGBA")
    draw_grid(draw, hex_tuple("#B8D9F5"), step=72, alpha=18)
    draw_glow(img, (1440, 210), 520, a1, 125)
    draw_glow(img, (1720, 850), 450, a2, 90)
    draw_glow(img, (280, 850), 360, hex_tuple("#0EA5E9"), 50)

    for i in range(13):
        y = 100 + i * 72
        draw.line([(1060, y), (1900, y + 220)], fill=a2 + (38,), width=2)

    for row in range(4):
        for col in range(5 - row):
            x = 1320 + col * 92 + row * 46
            y = 360 + row * 64
            draw.rounded_rectangle((x, y, x + 76, y + 44), radius=12, fill=(255, 255, 255, 20), outline=a2 + (72,), width=2)

    draw.rounded_rectangle((84, 88, 1836, 992), radius=44, outline=(255, 255, 255, 34), width=2)
    add_noise(img, alpha=5)
    return img.convert("RGB")


def ensure_backgrounds():
    ASSET_DIR.mkdir(exist_ok=True)
    paths = {}
    for key in ACCENTS:
        safe = section_slug(key)
        light = ASSET_DIR / f"bg_{safe}_light.png"
        dark = ASSET_DIR / f"bg_{safe}_dark.png"
        if not light.exists():
            create_light_background(key).save(light, optimize=True)
        if not dark.exists():
            create_dark_background(key).save(dark, optimize=True)
        paths[(key, "light")] = light
        paths[(key, "dark")] = dark
    return paths


def set_fill(shape, color, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency


def set_line(shape, color, width=1, transparency=None):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if transparency is not None:
        shape.line.transparency = transparency


def add_text(slide, text, x, y, w, h, size=24, color=None, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", fit=False):
    color = color or COLORS["ink"]
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if fit:
        tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_tag(slide, text, x, y, w, fill, fg=None, dark=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.36))
    set_fill(shape, fill, transparency=0 if dark else 8)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.045)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = fg or COLORS["white"]
    return shape


def is_dark_slide(idx, data):
    if idx in (1, 90):
        return True
    if idx in (4, 11, 18, 25, 40, 48, 58, 67, 77):
        return True
    if data["section"] == "2. Docker" and idx in (2, 3, 89):
        return False
    return False


def add_background(slide, data, idx, bg_paths, dark=False):
    key = section_key(data["section"])
    path = bg_paths[(key, "dark" if dark else "light")]
    slide.shapes.add_picture(str(path), 0, 0, Inches(13.333), Inches(7.5))
    accent = accent_rgb(data["section"])
    if not dark:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.48), Inches(0.42), Inches(12.36), Inches(6.62))
        set_fill(panel, rgb("FFFFFF"), transparency=10)
        set_line(panel, rgb("E2E8F0"), width=1, transparency=15)
    add_tag(
        slide,
        data["section"],
        Inches(0.76),
        Inches(0.54),
        Inches(2.72),
        accent,
        COLORS["white"],
        dark=True,
    )
    add_text(slide, f"{idx:02d}/90", Inches(11.95), Inches(0.6), Inches(0.62), Inches(0.2), 8,
             COLORS["white"] if dark else COLORS["muted"], bold=True, align=PP_ALIGN.RIGHT)


def add_progress(slide, idx, accent, dark=False):
    x, y, w = Inches(0.76), Inches(6.82), Inches(11.78)
    track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.035))
    set_fill(track, rgb("FFFFFF") if dark else rgb("D9E4EF"), transparency=70 if dark else 0)
    track.line.fill.background()
    done = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(11.78 * idx / 90), Inches(0.035))
    set_fill(done, accent)
    done.line.fill.background()


def add_content_title(slide, data, dark=False):
    title_color = COLORS["white"] if dark else COLORS["ink"]
    subtitle_color = rgb("D7E7F6") if dark else COLORS["muted"]
    size = 39 if dark else 32
    x = Inches(0.76)
    y = Inches(1.08 if dark else 1.0)
    w = Inches(7.25 if not dark else 8.0)
    h = Inches(1.22)
    add_text(slide, data["title"], x, y, w, h, size=size, color=title_color, bold=True, font="Aptos Display")
    add_text(slide, SECTION_LABELS[section_key(data["section"])], x, y + Inches(1.12), Inches(5.2), Inches(0.3), 12, subtitle_color, bold=True)


def add_bullet_cards(slide, bullets, accent, dark=False):
    x = Inches(0.82)
    y0 = Inches(2.68 if dark else 2.46)
    w = Inches(6.38 if not dark else 6.85)
    card_h = Inches(0.68)
    gap = Inches(0.19)
    fill = rgb("FFFFFF") if not dark else rgb("FFFFFF")
    fg = COLORS["ink"] if not dark else COLORS["white"]
    muted = COLORS["muted"] if not dark else rgb("D8E9F8")
    for i, bullet in enumerate(bullets):
        y = y0 + i * (card_h + gap)
        if dark:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, card_h)
            set_fill(card, rgb("FFFFFF"), transparency=88)
            set_line(card, rgb("FFFFFF"), width=1, transparency=78)
            text_color = fg
        else:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, card_h)
            set_fill(card, fill, transparency=0)
            set_line(card, rgb("E2E8F0"), width=1)
            text_color = COLORS["ink"]
        n = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.22), y + Inches(0.18), Inches(0.31), Inches(0.31))
        set_fill(n, accent)
        n.line.fill.background()
        add_text(slide, str(i + 1), x + Inches(0.22), y + Inches(0.26), Inches(0.31), Inches(0.08), 6, COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, bullet, x + Inches(0.68), y + Inches(0.18), w - Inches(0.9), Inches(0.32), 18, text_color, bold=False)
    if len(bullets) < 4:
        add_text(slide, "Cours Docker • débutant à intermédiaire", x, y0 + len(bullets) * (card_h + gap) + Inches(0.18),
                 Inches(4.8), Inches(0.25), 9, muted, bold=True)


def add_big_number(slide, idx, accent, dark=False):
    color = rgb("FFFFFF") if dark else rgb("E6EEF8")
    add_text(slide, f"{idx:02d}", Inches(8.15), Inches(0.96), Inches(3.2), Inches(1.2),
             64, color, bold=True, font="Aptos Display", align=PP_ALIGN.RIGHT)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.58), Inches(1.38), Inches(0.22), Inches(0.22))
    set_fill(dot, accent)
    dot.line.fill.background()


def add_terminal(slide, command, accent, dark=False):
    x, y, w, h = Inches(7.55), Inches(2.08), Inches(4.92), Inches(3.28)
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.08), y + Inches(0.1), w, h)
    set_fill(shadow, rgb("020617"), transparency=78 if not dark else 40)
    shadow.line.fill.background()

    term = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(term, COLORS["code"])
    set_line(term, rgb("1E293B"), width=1)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.5))
    set_fill(header, rgb("172033"))
    header.line.fill.background()
    for i, c in enumerate([COLORS["red"], COLORS["amber"], COLORS["green"]]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2 + i * 0.26), y + Inches(0.18), Inches(0.11), Inches(0.11))
        set_fill(dot, c)
        dot.line.fill.background()
    add_text(slide, "terminal", x + Inches(3.88), y + Inches(0.17), Inches(0.62), Inches(0.14), 7, rgb("94A3B8"), bold=True, align=PP_ALIGN.RIGHT)
    prompt = add_text(slide, "$", x + Inches(0.36), y + Inches(1.02), Inches(0.25), Inches(0.38), 23, accent, bold=True, font="Consolas")
    prompt.text_frame.paragraphs[0].runs[0].font.name = "Consolas"
    cmd_size = 20 if len(command) <= 30 else 17
    cmd = add_text(slide, command, x + Inches(0.72), y + Inches(1.02), w - Inches(1.05), Inches(0.9), cmd_size, rgb("F8FAFC"), font="Consolas")
    cmd.text_frame.paragraphs[0].runs[0].font.name = "Consolas"
    add_text(slide, "Une seule commande à retenir", x + Inches(0.36), y + Inches(2.42), Inches(3.5), Inches(0.24), 10, rgb("CBD5E1"), bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.36), y + Inches(2.16), Inches(1.42), Inches(0.035))
    set_fill(line, accent)
    line.line.fill.background()


def add_icon_badge(slide, accent, kind, dark=False):
    x, y = Inches(8.3), Inches(2.03)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(3.35), Inches(3.35))
    set_fill(badge, rgb("FFFFFF"), transparency=12 if not dark else 86)
    set_line(badge, rgb("FFFFFF") if dark else rgb("D9E4EF"), width=1, transparency=70 if dark else 0)
    cx, cy = x + Inches(1.675), y + Inches(1.675)

    if kind == "problem":
        for i, label in enumerate(["dev", "test", "prod"]):
            bx = x + Inches(0.58 + i * 0.76)
            by = y + Inches(1.04 + (i % 2) * 0.34)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(0.58), Inches(0.78))
            set_fill(box, accent if i == 2 else rgb("FFFFFF"), transparency=0 if i == 2 else 8)
            set_line(box, accent, width=1)
            add_text(slide, label, bx, by + Inches(0.3), Inches(0.58), Inches(0.11), 6,
                     COLORS["white"] if i == 2 else accent, bold=True, align=PP_ALIGN.CENTER)
        warn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x + Inches(1.2), y + Inches(2.15), Inches(0.9), Inches(0.78))
        set_fill(warn, accent)
        warn.line.fill.background()
        add_text(slide, "!", x + Inches(1.2), y + Inches(2.31), Inches(0.9), Inches(0.22), 18, COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    elif kind == "container":
        for i in range(3):
            cont = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.7 + i * 0.62), y + Inches(0.95), Inches(0.52), Inches(1.18))
            set_fill(cont, accent if i == 1 else rgb("FFFFFF"), transparency=0 if i == 1 else 10)
            set_line(cont, accent, width=1.2)
        base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.62), y + Inches(2.45), Inches(2.1), Inches(0.27))
        set_fill(base, accent)
        base.line.fill.background()
    elif kind == "compare":
        left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.62), y + Inches(0.74), Inches(0.9), Inches(2.02))
        right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(1.82), y + Inches(0.74), Inches(0.9), Inches(2.02))
        for shape in [left, right]:
            set_fill(shape, rgb("FFFFFF"), transparency=5)
            set_line(shape, accent, width=1)
        add_text(slide, "VM", x + Inches(0.62), y + Inches(1.55), Inches(0.9), Inches(0.18), 9, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "CT", x + Inches(1.82), y + Inches(1.55), Inches(0.9), Inches(0.18), 9, accent, bold=True, align=PP_ALIGN.CENTER)
    elif kind == "network":
        pts = [(1.67, .78), (.85, 1.7), (2.49, 1.7), (1.67, 2.55)]
        centers = [(x + Inches(px), y + Inches(py)) for px, py in pts]
        for a, b in [(0, 1), (0, 2), (1, 3), (2, 3), (1, 2)]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, centers[a][0], centers[a][1], centers[b][0], centers[b][1])
            line.line.color.rgb = accent
            line.line.width = Pt(1.4)
            line.line.transparency = 18
        for i, (nx, ny) in enumerate(centers):
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, nx - Inches(0.17), ny - Inches(0.17), Inches(0.34), Inches(0.34))
            set_fill(node, accent if i == 0 else rgb("FFFFFF"), transparency=0 if i == 0 else 8)
            set_line(node, accent, width=1)
    elif kind == "image":
        for i in range(5):
            layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.78 + i * 0.09), y + Inches(0.86 + i * 0.34), Inches(1.7 - i * 0.05), Inches(0.28))
            set_fill(layer, accent if i == 0 else rgb("FFFFFF"), transparency=0 if i == 0 else 7)
            set_line(layer, accent, width=1)
    elif kind == "hub":
        cloud = slide.shapes.add_shape(MSO_SHAPE.CLOUD, x + Inches(0.72), y + Inches(0.8), Inches(1.9), Inches(1.14))
        set_fill(cloud, rgb("FFFFFF"), transparency=5)
        set_line(cloud, accent, width=1.2)
        up = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, x + Inches(1.73), y + Inches(2.05), Inches(0.48), Inches(0.62))
        set_fill(up, accent)
        up.line.fill.background()
        down = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x + Inches(1.09), y + Inches(2.05), Inches(0.48), Inches(0.62))
        set_fill(down, accent)
        down.line.fill.background()
    elif kind == "storage":
        cyl = slide.shapes.add_shape(MSO_SHAPE.CAN, x + Inches(1.12), y + Inches(0.82), Inches(1.15), Inches(1.78))
        set_fill(cyl, rgb("FFFFFF"), transparency=4)
        set_line(cyl, accent, width=1.4)
        link = slide.shapes.add_shape(MSO_SHAPE.CUBE, x + Inches(0.7), y + Inches(1.45), Inches(0.62), Inches(0.62))
        set_fill(link, accent)
        link.line.fill.background()
    elif kind == "compose":
        centers = [(cx, cy - Inches(0.7)), (cx - Inches(0.72), cy + Inches(0.42)), (cx + Inches(0.72), cy + Inches(0.42))]
        for a, b in [(0, 1), (0, 2), (1, 2)]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, centers[a][0], centers[a][1], centers[b][0], centers[b][1])
            line.line.color.rgb = accent
            line.line.width = Pt(1.4)
        for i, (nx, ny) in enumerate(centers):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx - Inches(0.36), ny - Inches(0.24), Inches(0.72), Inches(0.48))
            set_fill(box, accent if i == 0 else rgb("FFFFFF"), transparency=0 if i == 0 else 7)
            set_line(box, accent, width=1)
    else:
        for i in range(3):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.72 + i * 0.33), y + Inches(0.95 + i * 0.43), Inches(1.55), Inches(0.9))
            set_fill(box, accent if i == 2 else rgb("FFFFFF"), transparency=0 if i == 2 else 8)
            set_line(box, accent, width=1)


def visual_kind(data):
    if data["command"]:
        return "command"
    mapping = {
        "problem": "problem",
        "container": "container",
        "compare": "compare",
        "install": "generic",
        "network": "network",
        "image": "image",
        "hub": "hub",
        "storage": "storage",
        "compose": "compose",
        "overview": "generic",
        "roadmap": "generic",
        "summary": "generic",
        "title": "generic",
    }
    return mapping.get(data.get("visual"), "generic")


def add_title_slide(slide, data, accent):
    add_text(slide, "FORMATION DEVOPS", Inches(0.86), Inches(0.76), Inches(2.8), Inches(0.28), 11, rgb("A5F3FC"), bold=True)
    add_text(slide, data["title"], Inches(0.84), Inches(1.42), Inches(7.5), Inches(1.75), 49, COLORS["white"], bold=True, font="Aptos Display")
    add_text(slide, "Concepts, utilisation et bonnes pratiques", Inches(0.88), Inches(3.18), Inches(5.9), Inches(0.42), 19, rgb("D7E7F6"))
    x0, y0 = Inches(0.9), Inches(4.35)
    for i, bullet in enumerate(data["bullets"]):
        add_tag(slide, bullet, x0 + Inches(i * 2.3), y0, Inches(2.05), accent, COLORS["white"], dark=True)


def add_slide(prs, data, idx, bg_paths):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark = is_dark_slide(idx, data)
    accent = accent_rgb(data["section"])
    add_background(slide, data, idx, bg_paths, dark=dark)

    if idx == 1:
        add_title_slide(slide, data, accent)
    else:
        add_big_number(slide, idx, accent, dark=dark)
        add_content_title(slide, data, dark=dark)
        add_bullet_cards(slide, data["bullets"], accent, dark=dark)
        if data["command"]:
            add_terminal(slide, data["command"], accent, dark=dark)
        else:
            add_icon_badge(slide, accent, visual_kind(data), dark=dark)

    add_progress(slide, idx, accent, dark=dark)
    return slide


def build_deck():
    assert len(SLIDES) == 90, f"Nombre de slides source incorrect : {len(SLIDES)}"
    bg_paths = ensure_backgrounds()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for idx, data in enumerate(SLIDES, start=1):
        add_slide(prs, data, idx, bg_paths)
    assert len(prs.slides) == 90, f"Nombre de slides incorrect : {len(prs.slides)}"
    props = prs.core_properties
    props.title = "Docker - Concepts, utilisation et bonnes pratiques"
    props.subject = "Présentation premium pour cours Docker"
    props.author = "OpenAI Codex"
    props.keywords = "Docker, DevOps, conteneurisation, Docker Compose, cours"
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build_deck()
    print(f"Présentation premium créée : {path.resolve()}")
    print(f"Nombre de slides : {len(SLIDES)}")
