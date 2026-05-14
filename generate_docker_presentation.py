# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("Docker_Concepts_Utilisation_Bonnes_Pratiques_90_slides.pptx")


def rgb(value):
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


COLORS = {
    "bg": rgb("F7FAFC"),
    "ink": rgb("172033"),
    "muted": rgb("5F6B7A"),
    "line": rgb("D8E0EA"),
    "panel": rgb("FFFFFF"),
    "navy": rgb("12304A"),
    "blue": rgb("1D63ED"),
    "cyan": rgb("1CA7B8"),
    "teal": rgb("168A7A"),
    "green": rgb("2F8F46"),
    "amber": rgb("E79A2F"),
    "orange": rgb("D96B2B"),
    "red": rgb("C94747"),
    "purple": rgb("7057C7"),
    "slate": rgb("4B6078"),
}


SECTION_COLORS = {
    "2.1": COLORS["orange"],
    "2.2": COLORS["teal"],
    "2.3": COLORS["blue"],
    "2.4": COLORS["navy"],
    "2.5": COLORS["cyan"],
    "2.6": COLORS["purple"],
    "2.7": COLORS["green"],
    "2.8": COLORS["amber"],
    "2.9": COLORS["slate"],
}


def S(section, title, bullets, visual="generic", command=None):
    return {
        "section": section,
        "title": title,
        "bullets": bullets,
        "visual": visual,
        "command": command,
    }


SLIDES = [
    S(
        "2. Docker",
        "Docker : concepts, utilisation et bonnes pratiques",
        ["Comprendre les conteneurs", "Lancer des applications", "Travailler proprement en équipe"],
        "title",
    ),
    S(
        "2. Docker",
        "Objectifs de la formation",
        ["Identifier les bons usages", "Manipuler images et conteneurs", "Préparer des projets reproductibles"],
        "roadmap",
    ),
    S(
        "2. Docker",
        "Docker en trois briques",
        ["Image : modèle de départ", "Conteneur : application lancée", "Engine : moteur d’exécution"],
        "overview",
    ),
    S(
        "2.1 Problématique (Pourquoi Docker ?)",
        "Problèmes avant Docker",
        ["Applications installées à la main", "Dépendances dispersées", "Corrections lentes et risquées"],
        "problem",
    ),
    S(
        "2.1 Problématique (Pourquoi Docker ?)",
        "Limites des installations classiques",
        ["Chaque serveur devient unique", "Mises à jour manuelles", "Retour arrière compliqué"],
        "problem",
    ),
    S(
        "2.1 Problématique (Pourquoi Docker ?)",
        "Environnements incohérents",
        ["Développement, test et production diffèrent", "Versions parfois incompatibles", "Erreurs difficiles à reproduire"],
        "problem",
    ),
    S(
        "2.1 Problématique (Pourquoi Docker ?)",
        "« It works on my machine »",
        ["Le poste local cache le vrai problème", "La production n’a pas les mêmes réglages", "Le diagnostic prend du temps"],
        "problem",
    ),
    S(
        "2.1 Problématique (Pourquoi Docker ?)",
        "Besoin de portabilité",
        ["Emballer l’application avec ses dépendances", "Déplacer le même paquet partout", "Réduire les surprises au déploiement"],
        "problem",
    ),
    S(
        "2.1 Problématique (Pourquoi Docker ?)",
        "Besoin de rapidité",
        ["Installer moins de choses à la main", "Démarrer en quelques secondes", "Tester plus souvent"],
        "problem",
    ),
    S(
        "2.1 Problématique (Pourquoi Docker ?)",
        "Besoin de standardisation",
        ["Même format pour livrer", "Mêmes commandes pour exécuter", "Pipeline CI/CD plus simple"],
        "problem",
    ),
    S(
        "2.2 Conteneurisation",
        "Définition de la conteneurisation",
        ["Isoler une application dans un conteneur", "Inclure ses dépendances utiles", "Partager le noyau de la machine hôte"],
        "container",
    ),
    S(
        "2.2 Conteneurisation",
        "Conteneur vs machine virtuelle",
        ["Une VM embarque un OS complet", "Un conteneur partage le noyau", "Le conteneur est souvent plus léger"],
        "compare",
    ),
    S(
        "2.2 Conteneurisation",
        "Isolation",
        ["Processus séparés", "Système de fichiers isolé", "Ressources contrôlables"],
        "container",
    ),
    S(
        "2.2 Conteneurisation",
        "Légèreté",
        ["Pas d’OS complet par application", "Images souvent petites", "Moins de mémoire utilisée"],
        "container",
    ),
    S(
        "2.2 Conteneurisation",
        "Démarrage rapide",
        ["Lancement en secondes", "Cycle test plus court", "Redémarrage simple après incident"],
        "container",
    ),
    S(
        "2.2 Conteneurisation",
        "Cas d’usage",
        ["API web", "Base de données de test", "Jobs batch et outils CLI", "Environnements de formation"],
        "container",
    ),
    S(
        "2.2 Conteneurisation",
        "Avantages clés",
        ["Portabilité", "Reproductibilité", "Déploiements plus prévisibles", "Meilleure isolation"],
        "container",
    ),
    S(
        "2.3 Installation de Docker",
        "Docker Engine",
        ["Moteur principal de Docker", "Gère images, conteneurs et réseaux", "Expose une API et une CLI"],
        "install",
    ),
    S(
        "2.3 Installation de Docker",
        "Docker Desktop",
        ["Application graphique pour poste de travail", "Inclut Docker Engine et Compose", "Pratique sur Windows et macOS"],
        "install",
    ),
    S(
        "2.3 Installation de Docker",
        "Installation Linux",
        ["Installer depuis le dépôt officiel", "Activer le service Docker", "Ajouter l’utilisateur au groupe si besoin"],
        "install",
    ),
    S(
        "2.3 Installation de Docker",
        "Installation Windows",
        ["Installer Docker Desktop", "Activer WSL 2 si demandé", "Redémarrer puis tester la CLI"],
        "install",
    ),
    S(
        "2.3 Installation de Docker",
        "Installation macOS",
        ["Installer Docker Desktop", "Choisir la version Intel ou Apple Silicon", "Lancer Docker avant la CLI"],
        "install",
    ),
    S(
        "2.3 Installation de Docker",
        "Vérification de l’installation",
        ["Contrôler que le service répond", "Lancer un conteneur de test", "Valider les droits utilisateur"],
        "install",
    ),
    S(
        "2.3 Installation de Docker",
        "Commande docker --version",
        ["Affiche la version installée", "Confirme que la CLI est disponible", "Premier test rapide après installation"],
        "command",
        "docker --version",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker pull télécharge une image",
        ["Récupère une image depuis un registre", "Prépare un futur lancement", "Utilise Docker Hub par défaut"],
        "command",
        "docker pull nginx",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker run lance un conteneur",
        ["Crée un conteneur depuis une image", "Démarre le processus principal", "Télécharge l’image si elle manque"],
        "command",
        "docker run nginx",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker run -it ouvre un mode interactif",
        ["Connecte le terminal au conteneur", "Utile pour explorer ou déboguer", "Souvent utilisé avec une image Linux"],
        "command",
        "docker run -it ubuntu",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker run -p publie un port",
        ["Relie un port local au conteneur", "Rend un service accessible", "Format courant : hôte:conteneur"],
        "command",
        "docker run -p 8080:80 nginx",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker run -d lance en arrière-plan",
        ["Détache le conteneur du terminal", "Retourne un identifiant", "Pratique pour les services longs"],
        "command",
        "docker run -d nginx",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker run -v monte un stockage",
        ["Connecte un volume ou un dossier", "Conserve les données importantes", "Sépare données et conteneur"],
        "command",
        "docker run -v data:/app/data nginx",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker ps liste les conteneurs actifs",
        ["Affiche les conteneurs en cours", "Montre ports, noms et statuts", "Très utile pour diagnostiquer"],
        "command",
        "docker ps",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker ps -a liste tout l’historique",
        ["Inclut les conteneurs arrêtés", "Aide à comprendre les échecs", "Permet de retrouver un ancien conteneur"],
        "command",
        "docker ps -a",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker inspect affiche les détails",
        ["Retourne la configuration complète", "Montre réseau, volumes et variables", "Sortie détaillée en JSON"],
        "command",
        "docker inspect <conteneur>",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker logs lit la sortie",
        ["Affiche les logs du conteneur", "Aide au diagnostic rapide", "Suit la sortie de l’application"],
        "command",
        "docker logs <conteneur>",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker exec exécute dans un conteneur",
        ["Lance une commande dans un conteneur actif", "Utile pour vérifier l’état interne", "À utiliser avec prudence en production"],
        "command",
        "docker exec -it <conteneur> sh",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker stop arrête proprement",
        ["Demande l’arrêt du processus principal", "Laisse le conteneur exister", "Évite de tuer brutalement le service"],
        "command",
        "docker stop <conteneur>",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker rm supprime un conteneur",
        ["Supprime un conteneur arrêté", "Libère de l’espace et de la liste", "Ne supprime pas l’image utilisée"],
        "command",
        "docker rm <conteneur>",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker images liste les images",
        ["Affiche les images locales", "Montre tags, tailles et identifiants", "Aide à nettoyer le poste"],
        "command",
        "docker images",
    ),
    S(
        "2.4 Commandes Docker de base",
        "docker rmi supprime une image",
        ["Supprime une image locale", "Possible si aucun conteneur ne l’utilise", "Libère de l’espace disque"],
        "command",
        "docker rmi <image>",
    ),
    S(
        "2.5 Réseaux sur Docker",
        "Réseau par défaut",
        ["Docker crée un réseau automatiquement", "Les conteneurs peuvent sortir vers Internet", "La configuration suffit pour les tests simples"],
        "network",
    ),
    S(
        "2.5 Réseaux sur Docker",
        "Driver bridge",
        ["Réseau local isolé sur l’hôte", "Mode par défaut le plus courant", "Idéal pour relier plusieurs conteneurs"],
        "network",
    ),
    S(
        "2.5 Réseaux sur Docker",
        "Driver host",
        ["Le conteneur partage le réseau de l’hôte", "Moins d’isolation réseau", "Utile pour certains besoins de performance"],
        "network",
    ),
    S(
        "2.5 Réseaux sur Docker",
        "Autres drivers",
        ["overlay pour plusieurs hôtes", "macvlan pour une IP dédiée", "none pour couper le réseau"],
        "network",
    ),
    S(
        "2.5 Réseaux sur Docker",
        "Réseau personnalisé",
        ["Nommer son réseau", "Mieux isoler les applications", "Faciliter la communication par nom"],
        "network",
    ),
    S(
        "2.5 Réseaux sur Docker",
        "Communication entre conteneurs",
        ["Les services partagent un réseau", "Ils s’appellent par nom", "Les ports internes restent simples"],
        "network",
    ),
    S(
        "2.5 Réseaux sur Docker",
        "DNS intégré Docker",
        ["Docker résout les noms de conteneurs", "Pas besoin d’IP fixes", "Les noms restent plus lisibles"],
        "network",
    ),
    S(
        "2.5 Réseaux sur Docker",
        "Cas pratiques réseau",
        ["API reliée à une base de données", "Proxy devant plusieurs services", "Environnement de test isolé"],
        "network",
    ),
    S(
        "2.6 Images Docker",
        "Qu’est-ce qu’une image Docker",
        ["Modèle immuable d’un conteneur", "Contient application et dépendances", "Se partage via un registre"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Structure en couches",
        ["Chaque étape ajoute une couche", "Les couches sont mises en cache", "Les couches communes sont réutilisées"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Dockerfile",
        ["Fichier texte de construction", "Décrit l’image étape par étape", "Rend la création reproductible"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Instructions Dockerfile essentielles",
        ["FROM choisit l’image de base", "RUN exécute une étape", "COPY ajoute des fichiers", "CMD définit le démarrage"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Processus de build",
        ["Docker lit le Dockerfile", "Construit les couches dans l’ordre", "Produit une image locale"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Docker builder",
        ["Moteur de construction des images", "Utilise le cache pour accélérer", "Supporte des builds plus avancés"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Multi-stage build",
        ["Sépare compilation et exécution", "Copie seulement le résultat utile", "Réduit fortement la taille finale"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Optimisation des images",
        ["Choisir une base adaptée", "Réduire les fichiers inutiles", "Regrouper les étapes cohérentes"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Taguer une image",
        ["Donne un nom lisible à l’image", "Ajoute une version claire", "Facilite push, pull et rollback"],
        "image",
    ),
    S(
        "2.6 Images Docker",
        "Bonnes pratiques images",
        ["Éviter les secrets dans l’image", "Utiliser des tags explicites", "Scanner les vulnérabilités", "Documenter le démarrage"],
        "image",
    ),
    S(
        "2.7 Docker Hub",
        "Qu’est-ce que Docker Hub",
        ["Plateforme de partage d’images", "Permet de publier et télécharger", "Très utilisée par la communauté"],
        "hub",
    ),
    S(
        "2.7 Docker Hub",
        "Registre public",
        ["Images accessibles depuis Internet", "Pratique pour les projets ouverts", "À éviter pour les contenus sensibles"],
        "hub",
    ),
    S(
        "2.7 Docker Hub",
        "Créer un compte Docker Hub",
        ["Compte nécessaire pour publier", "Identifiant utilisé dans les noms d’images", "Accès aux dépôts personnels"],
        "hub",
    ),
    S(
        "2.7 Docker Hub",
        "Créer un dépôt",
        ["Un dépôt regroupe les tags d’une image", "Nom stable pour l’équipe", "Visibilité publique ou privée selon besoin"],
        "hub",
    ),
    S(
        "2.7 Docker Hub",
        "docker login authentifie la CLI",
        ["Connecte le terminal au compte", "Autorise push et accès privé", "À protéger sur les postes partagés"],
        "command",
        "docker login",
    ),
    S(
        "2.7 Docker Hub",
        "docker push publie une image",
        ["Envoie l’image vers le registre", "Nécessite un tag complet", "Rend l’image disponible à l’équipe"],
        "command",
        "docker push user/app:1.0",
    ),
    S(
        "2.7 Docker Hub",
        "docker pull depuis Docker Hub",
        ["Télécharge une image publiée", "Utilise le tag demandé", "Permet un déploiement reproductible"],
        "command",
        "docker pull nginx:latest",
    ),
    S(
        "2.7 Docker Hub",
        "Images officielles",
        ["Maintenues par des sources reconnues", "Bon point de départ", "Lire la documentation avant usage"],
        "hub",
    ),
    S(
        "2.7 Docker Hub",
        "Bonnes pratiques de publication",
        ["Publier des tags versionnés", "Ne jamais pousser de secrets", "Ajouter une description utile", "Scanner avant diffusion"],
        "hub",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Problème de la persistance des données",
        ["Un conteneur peut être supprimé", "Les données internes peuvent disparaître", "Le stockage doit être séparé"],
        "storage",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Volumes Docker",
        ["Stockage géré par Docker", "Indépendant du cycle du conteneur", "Bon choix pour les données applicatives"],
        "storage",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Bind mounts",
        ["Montent un dossier de l’hôte", "Très utiles en développement", "Dépendent du chemin local"],
        "storage",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Différences volume vs bind mount",
        ["Volume : géré par Docker", "Bind mount : chemin de la machine", "Choisir selon portabilité et contrôle"],
        "storage",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Créer un volume",
        ["Prépare un espace persistant", "Nom simple à réutiliser", "Peut être partagé par plusieurs conteneurs"],
        "command",
        "docker volume create app_data",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Lister les volumes",
        ["Affiche les volumes existants", "Aide à vérifier les ressources", "Utile avant un nettoyage"],
        "command",
        "docker volume ls",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Utiliser un volume",
        ["Monte le volume dans le conteneur", "Garde les données après arrêt", "Sépare état et exécution"],
        "command",
        "docker run -v app_data:/data nginx",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Inspecter un volume",
        ["Affiche le point de montage", "Montre le driver utilisé", "Aide à diagnostiquer le stockage"],
        "command",
        "docker volume inspect app_data",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Supprimer un volume",
        ["Efface un stockage inutilisé", "À faire avec prudence", "Peut supprimer des données importantes"],
        "command",
        "docker volume rm app_data",
    ),
    S(
        "2.8 Stockage sur Docker",
        "Cas pratiques de stockage",
        ["Base de données locale", "Dossier de code en développement", "Cache partagé entre builds"],
        "storage",
    ),
    S(
        "2.9 Docker Compose",
        "Qu’est-ce que Docker Compose",
        ["Outil pour décrire plusieurs services", "Utilise un fichier YAML", "Pilote l’ensemble avec une seule CLI"],
        "compose",
    ),
    S(
        "2.9 Docker Compose",
        "Pourquoi Docker Compose",
        ["Évite les longues commandes répétées", "Regroupe configuration et dépendances", "Facilite le travail en équipe"],
        "compose",
    ),
    S(
        "2.9 Docker Compose",
        "Fichier compose.yaml",
        ["Fichier central du projet", "Décrit services, réseaux et volumes", "Versionnable avec le code"],
        "compose",
    ),
    S(
        "2.9 Docker Compose",
        "Structure du fichier compose",
        ["services pour les conteneurs", "networks pour la communication", "volumes pour les données"],
        "compose",
    ),
    S(
        "2.9 Docker Compose",
        "Services",
        ["Un service décrit un type de conteneur", "Peut définir image, build et ports", "Peut dépendre d’autres services"],
        "compose",
    ),
    S(
        "2.9 Docker Compose",
        "Networks dans Compose",
        ["Compose crée un réseau par projet", "Les services se joignent par nom", "Les réseaux peuvent être personnalisés"],
        "compose",
    ),
    S(
        "2.9 Docker Compose",
        "Volumes dans Compose",
        ["Déclarent les données persistantes", "Se montent dans les services", "Rendent l’environnement reproductible"],
        "compose",
    ),
    S(
        "2.9 Docker Compose",
        "docker compose build construit les images",
        ["Construit les services avec build", "Utilise les Dockerfile du projet", "Prépare un lancement cohérent"],
        "command",
        "docker compose build",
    ),
    S(
        "2.9 Docker Compose",
        "docker compose up démarre le projet",
        ["Lance tous les services nécessaires", "Crée réseaux et volumes manquants", "Affiche les logs par défaut"],
        "command",
        "docker compose up",
    ),
    S(
        "2.9 Docker Compose",
        "docker compose down arrête le projet",
        ["Arrête les services", "Supprime les conteneurs du projet", "Garde les volumes sauf option contraire"],
        "command",
        "docker compose down",
    ),
    S(
        "2.9 Docker Compose",
        "docker compose ps affiche les services",
        ["Liste l’état des conteneurs Compose", "Montre ports et noms", "Aide au diagnostic rapide"],
        "command",
        "docker compose ps",
    ),
    S(
        "2.9 Docker Compose",
        "docker compose logs lit les logs",
        ["Regroupe les sorties des services", "Permet de suivre un service précis", "Très utile pendant le debug"],
        "command",
        "docker compose logs",
    ),
    S(
        "2. Docker",
        "Synthèse des bonnes pratiques",
        ["Images petites et versionnées", "Données stockées hors conteneur", "Commandes automatisées avec Compose", "Secrets jamais dans les images"],
        "summary",
    ),
    S(
        "2. Docker",
        "Conclusion",
        ["Docker standardise l’exécution", "Les conteneurs rendent les tests rapides", "Compose simplifie les applications multi-services"],
        "summary",
    ),
]


def section_color(section):
    for prefix, color in SECTION_COLORS.items():
        if section.startswith(prefix):
            return color
    return COLORS["blue"]


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def set_line(shape, color=COLORS["line"], width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(slide, text, x, y, w, h, size=24, color=COLORS["ink"], bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", all_caps=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text.upper() if all_caps else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_label(slide, text, x, y, w, h, fill, color=RGBColor(255, 255, 255), size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(shape, fill)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return shape


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.82), Inches(6.7), Inches(4.1))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Aptos"
        p.font.size = Pt(23)
        p.font.color.rgb = COLORS["ink"]
        p.space_after = Pt(13)
        p.line_spacing = 1.08
    return box


def add_background(slide, data, idx):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    set_fill(bg, COLORS["bg"])
    bg.line.fill.background()

    accent = section_color(data["section"])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), Inches(7.5))
    set_fill(bar, accent)
    bar.line.fill.background()

    add_label(slide, data["section"], Inches(0.65), Inches(0.35), Inches(3.55), Inches(0.33), accent, size=9)
    add_text(slide, f"{idx:02d}/90", Inches(11.85), Inches(7.05), Inches(0.8), Inches(0.25), 9, COLORS["muted"], align=PP_ALIGN.RIGHT)

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(6.85), Inches(11.8), Inches(0.01))
    set_fill(rule, COLORS["line"])
    rule.line.fill.background()


def add_title_and_bullets(slide, data):
    add_text(
        slide,
        data["title"],
        Inches(0.65),
        Inches(0.78),
        Inches(7.5),
        Inches(0.82),
        size=30,
        color=COLORS["ink"],
        bold=True,
    )
    add_bullets(slide, data["bullets"])


def draw_terminal(slide, command, accent):
    x, y, w, h = Inches(8.0), Inches(1.72), Inches(4.55), Inches(2.65)
    term = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(term, COLORS["navy"])
    term.line.color.rgb = rgb("0C2237")
    term.line.width = Pt(1)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.42))
    set_fill(header, rgb("20364F"))
    header.line.fill.background()
    for i, c in enumerate([COLORS["red"], COLORS["amber"], COLORS["green"]]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.18 + i * 0.23), y + Inches(0.14), Inches(0.1), Inches(0.1)
        )
        set_fill(dot, c)
        dot.line.fill.background()

    prompt = add_text(slide, "$", x + Inches(0.32), y + Inches(0.95), Inches(0.3), Inches(0.35), 22, accent, bold=True, font="Consolas")
    prompt.text_frame.paragraphs[0].runs[0].font.name = "Consolas"
    cmd = add_text(slide, command, x + Inches(0.72), y + Inches(0.95), w - Inches(1.0), Inches(0.8), 20, RGBColor(255, 255, 255), font="Consolas")
    cmd.text_frame.paragraphs[0].runs[0].font.name = "Consolas"
    add_text(slide, "Commande clé", x + Inches(0.32), y + Inches(2.05), w - Inches(0.64), Inches(0.35), 11, rgb("B9C8DA"))


def draw_layer_stack(slide, accent, x=Inches(8.45), y=Inches(1.35)):
    labels = ["app", "deps", "runtime", "base"]
    widths = [3.2, 3.45, 3.7, 3.95]
    for i, label in enumerate(labels):
        layer = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches((4.0 - widths[i]) / 2),
            y + Inches(i * 0.72),
            Inches(widths[i]),
            Inches(0.52),
        )
        fill = accent if i == 0 else rgb(["D9EEF2", "E7EAF8", "EAF0F7"][min(i - 1, 2)])
        set_fill(layer, fill)
        set_line(layer, rgb("C5D2E0"))
        add_text(slide, label, x + Inches((4.0 - widths[i]) / 2), y + Inches(i * 0.72 + 0.13), Inches(widths[i]), Inches(0.25), 10, RGBColor(255, 255, 255) if i == 0 else COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)


def draw_container(slide, accent):
    x, y = Inches(8.15), Inches(1.35)
    host = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.35), Inches(3.75))
    set_fill(host, rgb("FFFFFF"))
    set_line(host, rgb("CAD6E2"), 1.2)
    add_text(slide, "HÔTE", x + Inches(0.22), y + Inches(0.22), Inches(0.8), Inches(0.25), 9, COLORS["muted"], bold=True)

    for i, label in enumerate(["APP", "DB"]):
        cx = x + Inches(0.42 + i * 2.05)
        cont = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y + Inches(0.82), Inches(1.55), Inches(1.9))
        set_fill(cont, rgb("EFF7F6") if i == 0 else rgb("F4F1FB"))
        set_line(cont, accent if i == 0 else COLORS["purple"], 1.1)
        add_text(slide, label, cx, y + Inches(1.02), Inches(1.55), Inches(0.3), 14, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "deps", cx, y + Inches(1.52), Inches(1.55), Inches(0.25), 10, COLORS["muted"], align=PP_ALIGN.CENTER)
        add_text(slide, "bin", cx, y + Inches(1.9), Inches(1.55), Inches(0.25), 10, COLORS["muted"], align=PP_ALIGN.CENTER)

    kernel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.42), y + Inches(3.0), Inches(3.55), Inches(0.38))
    set_fill(kernel, accent)
    kernel.line.fill.background()
    add_text(slide, "noyau partagé", x + Inches(0.42), y + Inches(3.08), Inches(3.55), Inches(0.18), 9, RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)


def draw_problem(slide, accent):
    x, y = Inches(8.05), Inches(1.35)
    names = ["DEV", "TEST", "PROD"]
    fills = [rgb("FFF3E9"), rgb("FFF9DF"), rgb("FEEDEE")]
    for i, name in enumerate(names):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(i * 1.45),
            y + Inches(0.4 + (i % 2) * 0.35),
            Inches(1.18),
            Inches(1.4),
        )
        set_fill(box, fills[i])
        set_line(box, accent if i == 2 else rgb("D8B99A"))
        add_text(slide, name, x + Inches(i * 1.45), y + Inches(0.65 + (i % 2) * 0.35), Inches(1.18), Inches(0.25), 10, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, f"v{i + 1}", x + Inches(i * 1.45), y + Inches(1.15 + (i % 2) * 0.35), Inches(1.18), Inches(0.22), 10, COLORS["muted"], align=PP_ALIGN.CENTER)

    warn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x + Inches(1.28), y + Inches(2.65), Inches(1.3), Inches(1.15))
    set_fill(warn, accent)
    warn.line.fill.background()
    add_text(slide, "!", x + Inches(1.28), y + Inches(2.88), Inches(1.3), Inches(0.4), 28, RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "écarts", x + Inches(0.7), y + Inches(4.05), Inches(2.5), Inches(0.3), 12, COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_compare(slide, accent):
    x, y = Inches(8.0), Inches(1.3)
    for i, title in enumerate(["VM", "CONTENEUR"]):
        bx = x + Inches(i * 2.15)
        outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, Inches(1.85), Inches(3.65))
        set_fill(outer, rgb("FFFFFF"))
        set_line(outer, rgb("CAD6E2"))
        add_text(slide, title, bx, y + Inches(0.25), Inches(1.85), Inches(0.3), 11, accent if i else COLORS["slate"], bold=True, align=PP_ALIGN.CENTER)
        parts = ["APP", "OS", "HYP"] if i == 0 else ["APP", "DEPS", "NOYAU"]
        for j, p in enumerate(parts):
            r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx + Inches(0.25), y + Inches(0.85 + j * 0.68), Inches(1.35), Inches(0.42))
            set_fill(r, accent if (i == 1 and j == 0) else rgb("EAF0F7"))
            r.line.fill.background()
            add_text(slide, p, bx + Inches(0.25), y + Inches(0.95 + j * 0.68), Inches(1.35), Inches(0.18), 8, RGBColor(255, 255, 255) if (i == 1 and j == 0) else COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)


def draw_install(slide, accent):
    x, y = Inches(8.15), Inches(1.35)
    screen = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.2), Inches(2.55))
    set_fill(screen, rgb("FFFFFF"))
    set_line(screen, rgb("C7D3DF"), 1.2)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.2), Inches(0.35))
    set_fill(header, accent)
    header.line.fill.background()
    add_text(slide, "Docker", x + Inches(0.22), y + Inches(0.08), Inches(1.1), Inches(0.18), 8, RGBColor(255, 255, 255), bold=True)
    add_text(slide, "CLI", x + Inches(0.45), y + Inches(0.9), Inches(0.9), Inches(0.35), 16, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "+", x + Inches(1.55), y + Inches(0.9), Inches(0.3), Inches(0.35), 18, accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Engine", x + Inches(2.0), y + Inches(0.9), Inches(1.3), Inches(0.35), 16, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    stand = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.75), y + Inches(2.55), Inches(0.7), Inches(0.28))
    set_fill(stand, rgb("C7D3DF"))
    stand.line.fill.background()
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(1.2), y + Inches(2.82), Inches(1.8), Inches(0.18))
    set_fill(base, rgb("C7D3DF"))
    base.line.fill.background()


def draw_network(slide, accent):
    x, y = Inches(8.25), Inches(1.35)
    points = [(1.9, 0.45), (0.65, 1.65), (3.1, 1.65), (1.9, 2.85)]
    centers = []
    for px, py in points:
        centers.append((x + Inches(px), y + Inches(py)))
    for a, b in [(0, 1), (0, 2), (1, 3), (2, 3), (1, 2)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, centers[a][0], centers[a][1], centers[b][0], centers[b][1])
        line.line.color.rgb = rgb("B7C9D8")
        line.line.width = Pt(1.3)
    for i, (cx, cy) in enumerate(centers):
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.32), cy - Inches(0.32), Inches(0.64), Inches(0.64))
        set_fill(node, accent if i == 0 else rgb("FFFFFF"))
        set_line(node, accent, 1.3)
        add_text(slide, str(i + 1), cx - Inches(0.32), cy - Inches(0.12), Inches(0.64), Inches(0.22), 9, RGBColor(255, 255, 255) if i == 0 else accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "réseau", x + Inches(0.72), y + Inches(3.55), Inches(2.5), Inches(0.3), 12, COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_hub(slide, accent):
    x, y = Inches(8.3), Inches(1.35)
    cloud = slide.shapes.add_shape(MSO_SHAPE.CLOUD, x + Inches(0.62), y + Inches(0.55), Inches(3.0), Inches(1.65))
    set_fill(cloud, rgb("FFFFFF"))
    set_line(cloud, accent, 1.4)
    add_text(slide, "REGISTRY", x + Inches(0.85), y + Inches(1.18), Inches(2.5), Inches(0.25), 11, accent, bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(["pull", "push"]):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.25 + i * 2.35), y + Inches(2.9), Inches(1.45), Inches(0.55))
        set_fill(box, accent if i else rgb("EAF7EF"))
        set_line(box, accent)
        add_text(slide, label, x + Inches(0.25 + i * 2.35), y + Inches(3.08), Inches(1.45), Inches(0.18), 9, RGBColor(255, 255, 255) if i else accent, bold=True, align=PP_ALIGN.CENTER)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(1.0 + i * 2.35), y + Inches(2.85), x + Inches(2.1), y + Inches(2.2))
        line.line.color.rgb = rgb("B7C9D8")
        line.line.width = Pt(1.2)


def draw_storage(slide, accent):
    x, y = Inches(8.45), Inches(1.35)
    for i, label in enumerate(["CONTENEUR", "VOLUME"]):
        bx = x + Inches(0.1 + i * 2.05)
        if i == 0:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y + Inches(0.65), Inches(1.55), Inches(1.3))
            set_fill(shape, rgb("FFFFFF"))
            set_line(shape, rgb("C8D4DF"))
        else:
            body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, y + Inches(0.65), Inches(1.55), Inches(1.3))
            set_fill(body, rgb("FFF7E8"))
            set_line(body, accent)
            top = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, y + Inches(0.48), Inches(1.55), Inches(0.35))
            set_fill(top, rgb("FFF7E8"))
            set_line(top, accent)
            bot = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, y + Inches(1.78), Inches(1.55), Inches(0.35))
            set_fill(bot, rgb("FFF7E8"))
            set_line(bot, accent)
        add_text(slide, label, bx, y + Inches(1.08), Inches(1.55), Inches(0.24), 9, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(1.66), y + Inches(1.3), x + Inches(2.08), y + Inches(1.3))
    line.line.color.rgb = accent
    line.line.width = Pt(2)
    add_text(slide, "persistant", x + Inches(0.75), y + Inches(3.2), Inches(2.6), Inches(0.3), 12, COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_compose(slide, accent):
    x, y = Inches(8.15), Inches(1.25)
    services = [("web", 1.25, 0.3), ("api", 0.3, 1.65), ("db", 2.2, 1.65)]
    centers = {}
    for label, px, py in services:
        bx, by = x + Inches(px), y + Inches(py)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(1.25), Inches(0.74))
        set_fill(box, accent if label == "web" else rgb("FFFFFF"))
        set_line(box, accent)
        add_text(slide, label, bx, by + Inches(0.23), Inches(1.25), Inches(0.2), 10, RGBColor(255, 255, 255) if label == "web" else accent, bold=True, align=PP_ALIGN.CENTER)
        centers[label] = (bx + Inches(0.625), by + Inches(0.37))
    for a, b in [("web", "api"), ("web", "db"), ("api", "db")]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, centers[a][0], centers[a][1], centers[b][0], centers[b][1])
        line.line.color.rgb = rgb("B7C9D8")
        line.line.width = Pt(1.2)
    file = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.72), y + Inches(3.0), Inches(2.7), Inches(0.62))
    set_fill(file, rgb("EEF2F8"))
    set_line(file, rgb("CAD6E2"))
    add_text(slide, "compose.yaml", x + Inches(0.72), y + Inches(3.22), Inches(2.7), Inches(0.18), 10, COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)


def draw_overview(slide, accent):
    x, y = Inches(8.05), Inches(1.45)
    labels = ["IMAGE", "CONTENEUR", "ENGINE"]
    for i, label in enumerate(labels):
        bx = x + Inches(0.2 + i * 1.35)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y + Inches(0.85), Inches(1.12), Inches(0.85))
        set_fill(shape, accent if i == 1 else rgb("FFFFFF"))
        set_line(shape, accent)
        add_text(slide, label, bx, y + Inches(1.15), Inches(1.12), Inches(0.18), 8, RGBColor(255, 255, 255) if i == 1 else accent, bold=True, align=PP_ALIGN.CENTER)
        if i < 2:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, bx + Inches(1.13), y + Inches(1.28), bx + Inches(1.35), y + Inches(1.28))
            line.line.color.rgb = rgb("B7C9D8")
            line.line.width = Pt(1.5)
    add_text(slide, "modèle → exécution → moteur", x + Inches(0.6), y + Inches(2.7), Inches(3.1), Inches(0.3), 11, COLORS["muted"], align=PP_ALIGN.CENTER)


def draw_roadmap(slide, accent):
    x, y = Inches(8.25), Inches(1.25)
    items = ["Comprendre", "Lancer", "Organiser"]
    for i, item in enumerate(items):
        cy = y + Inches(0.55 + i * 1.0)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), cy, Inches(0.44), Inches(0.44))
        set_fill(circle, accent)
        circle.line.fill.background()
        add_text(slide, str(i + 1), x + Inches(0.15), cy + Inches(0.12), Inches(0.44), Inches(0.15), 8, RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, item, x + Inches(0.85), cy + Inches(0.1), Inches(2.7), Inches(0.24), 15, COLORS["ink"], bold=True)
        if i < 2:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(0.37), cy + Inches(0.44), x + Inches(0.37), cy + Inches(1.0))
            line.line.color.rgb = rgb("B7C9D8")
            line.line.width = Pt(1.2)


def draw_summary(slide, accent):
    x, y = Inches(8.25), Inches(1.35)
    for i in range(4):
        row = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(i * 0.78), Inches(3.8), Inches(0.48))
        set_fill(row, rgb("FFFFFF"))
        set_line(row, rgb("D3DDE8"))
        mark = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.28), y + Inches(0.11 + i * 0.78), Inches(0.26), Inches(0.26))
        set_fill(mark, accent)
        mark.line.fill.background()
        add_text(slide, "✓", x + Inches(0.28), y + Inches(0.125 + i * 0.78), Inches(0.26), Inches(0.1), 9, RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)


def draw_title_visual(slide, accent):
    x, y = Inches(7.65), Inches(1.0)
    for row in range(3):
        for col in range(4 - row):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x + Inches(col * 0.82 + row * 0.41),
                y + Inches(row * 0.62),
                Inches(0.68),
                Inches(0.42),
            )
            set_fill(box, accent if row == 0 else rgb("DDEBFF"))
            box.line.fill.background()
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x - Inches(0.08), y + Inches(2.18), Inches(3.55), Inches(0.42))
    set_fill(base, COLORS["navy"])
    base.line.fill.background()
    add_text(slide, "containers", x, y + Inches(2.3), Inches(3.35), Inches(0.14), 8, RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)


def draw_generic(slide, accent):
    x, y = Inches(8.45), Inches(1.55)
    for i in range(3):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(i * 0.48), y + Inches(i * 0.48), Inches(2.7), Inches(1.4))
        set_fill(shape, rgb(["FFFFFF", "EEF4FA", "E7F3F0"][i]))
        set_line(shape, accent if i == 2 else rgb("CAD6E2"))


VISUALS = {
    "problem": draw_problem,
    "container": draw_container,
    "compare": draw_compare,
    "install": draw_install,
    "network": draw_network,
    "image": lambda slide, accent: draw_layer_stack(slide, accent),
    "hub": draw_hub,
    "storage": draw_storage,
    "compose": draw_compose,
    "overview": draw_overview,
    "roadmap": draw_roadmap,
    "summary": draw_summary,
    "title": draw_title_visual,
    "generic": draw_generic,
}


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        accent = section_color(data["section"])
        add_background(slide, data, idx)
        add_title_and_bullets(slide, data)

        if data["command"]:
            draw_terminal(slide, data["command"], accent)
        else:
            VISUALS.get(data["visual"], draw_generic)(slide, accent)

    assert len(prs.slides) == 90, f"Nombre de slides incorrect : {len(prs.slides)}"

    props = prs.core_properties
    props.title = "Docker - Concepts, utilisation et bonnes pratiques"
    props.subject = "Formation Docker débutant à intermédiaire"
    props.author = "OpenAI Codex"
    props.keywords = "Docker, DevOps, conteneurisation, Docker Compose"

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build_deck()
    print(f"Présentation créée : {path.resolve()}")
    print(f"Nombre de slides : {len(SLIDES)}")
