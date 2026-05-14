# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from generate_docker_presentation import SLIDES


OUT = Path("Docker_Cours_Haut_De_Gamme_Design_Unique_90_slides.pptx")


def rgb(value):
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


ACCENT = rgb("0B63CE")
ACCENT_DARK = rgb("074EA6")
ACCENT_LIGHT = rgb("DCEBFF")
ACCENT_PALE = rgb("F0F6FF")
BG = rgb("F8FAFD")
WHITE = rgb("FFFFFF")
INK = rgb("0F172A")
MUTED = rgb("64748B")
SOFT = rgb("E5EDF7")
LINE = rgb("D8E4F0")
NEUTRAL = rgb("EEF3F8")


SECTION_SHORT = {
    "2. Docker": "Vue d’ensemble",
    "2.1": "Pourquoi Docker ?",
    "2.2": "Conteneurisation",
    "2.3": "Installation",
    "2.4": "Commandes Docker",
    "2.5": "Réseaux",
    "2.6": "Images",
    "2.7": "Docker Hub",
    "2.8": "Stockage",
    "2.9": "Docker Compose",
}


def section_key(section):
    for key in ["2.1", "2.2", "2.3", "2.4", "2.5", "2.6", "2.7", "2.8", "2.9"]:
        if section.startswith(key):
            return key
    return "2. Docker"


def section_label(section):
    return SECTION_SHORT[section_key(section)]


def set_fill(shape, color, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1, transparency=None):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if transparency is not None:
        shape.line.transparency = transparency


def no_line(shape):
    shape.line.fill.background()


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_pill(slide, text, x, y, w, fill=ACCENT, color=WHITE):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34))
    set_fill(pill, fill)
    no_line(pill)
    tf = pill.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.035)
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = color
    return pill


def add_background(slide, data, idx):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(bg, BG)
    no_line(bg)

    # Bande claire, très sobre : elle remplace les fonds "datacenter".
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.55), 0, Inches(5.78), Inches(7.5))
    set_fill(band, ACCENT_PALE)
    no_line(band)

    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.075))
    set_fill(top_line, ACCENT)
    no_line(top_line)

    add_pill(slide, section_label(data["section"]), Inches(0.72), Inches(0.48), Inches(2.18), ACCENT, WHITE)
    add_text(slide, data["section"], Inches(3.08), Inches(0.54), Inches(4.0), Inches(0.2), 8.5, MUTED, bold=True)

    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(6.82), Inches(11.9), Inches(0.012))
    set_fill(footer, LINE)
    no_line(footer)
    progress = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(6.82), Inches(11.9 * idx / 90), Inches(0.034))
    set_fill(progress, ACCENT)
    no_line(progress)

    add_text(slide, "Docker – Concepts, utilisation et bonnes pratiques", Inches(0.72), Inches(7.02), Inches(5.6), Inches(0.18), 8, MUTED, bold=True)
    add_text(slide, f"{idx:02d} / 90", Inches(11.56), Inches(6.94), Inches(1.06), Inches(0.28), 12, ACCENT, bold=True, align=PP_ALIGN.RIGHT)


def add_title(slide, data, idx):
    is_title = idx == 1
    size = 42 if is_title else 34
    y = Inches(1.3 if is_title else 1.12)
    w = Inches(6.8 if not is_title else 7.3)
    add_text(slide, data["title"], Inches(0.72), y, w, Inches(1.24), size, INK, bold=True, font="Aptos Display")
    if is_title:
        add_text(slide, "Support de cours progressif • débutant à intermédiaire", Inches(0.76), Inches(2.82), Inches(6.1), Inches(0.34), 17, MUTED)


def add_bullets(slide, bullets, idx):
    if idx == 1:
        x, y0, w = Inches(0.78), Inches(4.15), Inches(6.2)
        for i, bullet in enumerate(bullets):
            add_pill(slide, bullet, x + Inches(i * 2.05), y0, Inches(1.82), ACCENT_LIGHT, ACCENT_DARK)
        return

    x, y0, w = Inches(0.78), Inches(2.76), Inches(6.25)
    card_h = Inches(0.72)
    gap = Inches(0.18)
    for i, bullet in enumerate(bullets):
        y = y0 + i * (card_h + gap)
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.035), y + Inches(0.045), w, card_h)
        set_fill(shadow, rgb("CBD5E1"), transparency=82)
        no_line(shadow)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, card_h)
        set_fill(card, WHITE)
        set_line(card, LINE, width=0.9)

        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.24), y + Inches(0.22), Inches(0.27), Inches(0.27))
        set_fill(dot, ACCENT)
        no_line(dot)
        add_text(slide, str(i + 1), x + Inches(0.24), y + Inches(0.285), Inches(0.27), Inches(0.08), 5.5, WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, bullet, x + Inches(0.68), y + Inches(0.2), w - Inches(0.92), Inches(0.3), 18.5, INK)


def add_visual_panel(slide):
    x, y, w, h = Inches(7.72), Inches(1.32), Inches(4.72), Inches(4.95)
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.05), y + Inches(0.07), w, h)
    set_fill(shadow, rgb("C7D2E1"), transparency=72)
    no_line(shadow)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(panel, WHITE)
    set_line(panel, LINE, width=1)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.28), y + Inches(0.28), Inches(0.64), Inches(0.12))
    set_fill(accent_bar, ACCENT)
    no_line(accent_bar)
    return x, y, w, h


def add_visual_title(slide, x, y, text):
    add_text(slide, text.upper(), x + Inches(0.28), y + Inches(0.54), Inches(3.7), Inches(0.22), 9, MUTED, bold=True)


def draw_title_visual(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Architecture moderne")
    bx, by = x + Inches(0.75), y + Inches(1.36)
    for row in range(3):
        for col in range(4 - row):
            c = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bx + Inches(col * 0.68 + row * 0.34),
                by + Inches(row * 0.54),
                Inches(0.54),
                Inches(0.34),
            )
            set_fill(c, ACCENT if row == 0 else ACCENT_LIGHT)
            no_line(c)
    ship = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx - Inches(0.14), by + Inches(2.06), Inches(3.25), Inches(0.56))
    set_fill(ship, ACCENT_DARK)
    no_line(ship)
    add_text(slide, "DOCKER", bx, by + Inches(2.24), Inches(2.98), Inches(0.12), 8, WHITE, bold=True, align=PP_ALIGN.CENTER)


def draw_generic(slide, label="Concept clé"):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, label)
    cx, cy = x + Inches(2.36), y + Inches(2.65)
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(1.1), cy - Inches(1.1), Inches(2.2), Inches(2.2))
    set_fill(outer, ACCENT_PALE)
    set_line(outer, ACCENT_LIGHT, width=2)
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.62), cy - Inches(0.62), Inches(1.24), Inches(1.24))
    set_fill(inner, ACCENT)
    no_line(inner)
    add_text(slide, "✓", cx - Inches(0.62), cy - Inches(0.28), Inches(1.24), Inches(0.32), 26, WHITE, bold=True, align=PP_ALIGN.CENTER)


def draw_problem(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Avant Docker")
    labels = ["DEV", "TEST", "PROD"]
    for i, label in enumerate(labels):
        bx = x + Inches(0.54 + i * 1.25)
        by = y + Inches(1.65 + (i % 2) * 0.35)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(0.92), Inches(1.22))
        set_fill(box, ACCENT if i == 2 else ACCENT_PALE)
        set_line(box, ACCENT_LIGHT if i < 2 else ACCENT, width=1.2)
        add_text(slide, label, bx, by + Inches(0.42), Inches(0.92), Inches(0.16), 9, WHITE if i == 2 else ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(0.96), y + Inches(3.35), x + Inches(3.72), y + Inches(3.35))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2.2)
    add_text(slide, "environnements différents", x + Inches(0.72), y + Inches(3.62), Inches(3.2), Inches(0.22), 10, MUTED, bold=True, align=PP_ALIGN.CENTER)


def draw_container(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Conteneurs")
    host = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.7), y + Inches(1.2), Inches(3.25), Inches(3.1))
    set_fill(host, ACCENT_PALE)
    set_line(host, ACCENT_LIGHT, width=1.4)
    for i, label in enumerate(["app", "api", "db"]):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(1.0 + i * 0.82), y + Inches(1.8), Inches(0.64), Inches(1.2))
        set_fill(box, WHITE)
        set_line(box, ACCENT, width=1.1)
        add_text(slide, label, x + Inches(1.0 + i * 0.82), y + Inches(2.3), Inches(0.64), Inches(0.14), 8, ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(1.0), y + Inches(3.45), Inches(2.3), Inches(0.32))
    set_fill(base, ACCENT)
    no_line(base)
    add_text(slide, "noyau partagé", x + Inches(1.0), y + Inches(3.55), Inches(2.3), Inches(0.1), 7, WHITE, bold=True, align=PP_ALIGN.CENTER)


def draw_compare(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Comparer")
    for i, title in enumerate(["VM", "Conteneur"]):
        bx = x + Inches(0.68 + i * 2.0)
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y + Inches(1.24), Inches(1.44), Inches(2.86))
        set_fill(frame, ACCENT_PALE if i == 0 else WHITE)
        set_line(frame, ACCENT_LIGHT if i == 0 else ACCENT, width=1.2)
        add_text(slide, title, bx, y + Inches(1.52), Inches(1.44), Inches(0.2), 10, ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)
        blocks = 4 if i == 0 else 2
        for j in range(blocks):
            block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx + Inches(0.28), y + Inches(2.0 + j * 0.36), Inches(0.88), Inches(0.22))
            set_fill(block, ACCENT if j == 0 else ACCENT_LIGHT)
            no_line(block)


def draw_network(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Réseau")
    pts = [(2.36, 1.42), (1.25, 2.45), (3.47, 2.45), (2.36, 3.5)]
    centers = [(x + Inches(px), y + Inches(py)) for px, py in pts]
    for a, b in [(0, 1), (0, 2), (1, 3), (2, 3), (1, 2)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, centers[a][0], centers[a][1], centers[b][0], centers[b][1])
        line.line.color.rgb = ACCENT_LIGHT
        line.line.width = Pt(2)
    for i, (cx, cy) in enumerate(centers):
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.25), cy - Inches(0.25), Inches(0.5), Inches(0.5))
        set_fill(node, ACCENT if i == 0 else WHITE)
        set_line(node, ACCENT, width=1.4)


def draw_image(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Image Docker")
    widths = [3.1, 2.82, 2.54, 2.26, 1.98]
    labels = ["app", "deps", "runtime", "os", "base"]
    for i, width in enumerate(widths):
        bx = x + Inches((4.72 - width) / 2)
        by = y + Inches(1.42 + i * 0.48)
        layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(width), Inches(0.34))
        set_fill(layer, ACCENT if i == 0 else ACCENT_PALE)
        set_line(layer, ACCENT_LIGHT, width=1)
        add_text(slide, labels[i], bx, by + Inches(0.095), Inches(width), Inches(0.08), 7, WHITE if i == 0 else ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)


def draw_hub(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Registre")
    cloud = slide.shapes.add_shape(MSO_SHAPE.CLOUD, x + Inches(1.22), y + Inches(1.28), Inches(2.22), Inches(1.25))
    set_fill(cloud, ACCENT_PALE)
    set_line(cloud, ACCENT, width=1.4)
    add_text(slide, "HUB", x + Inches(1.44), y + Inches(1.78), Inches(1.75), Inches(0.18), 13, ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(["pull", "push"]):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW if i == 0 else MSO_SHAPE.LEFT_ARROW,
                                       x + Inches(0.9 + i * 2.25), y + Inches(3.1), Inches(0.9), Inches(0.34))
        set_fill(arrow, ACCENT)
        no_line(arrow)
        add_text(slide, label, x + Inches(0.7 + i * 2.22), y + Inches(3.58), Inches(1.28), Inches(0.14), 8, MUTED, bold=True, align=PP_ALIGN.CENTER)


def draw_storage(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Persistance")
    cont = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.78), y + Inches(1.65), Inches(1.3), Inches(1.3))
    set_fill(cont, WHITE)
    set_line(cont, ACCENT, width=1.3)
    add_text(slide, "conteneur", x + Inches(0.78), y + Inches(2.15), Inches(1.3), Inches(0.14), 7, ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)
    cyl = slide.shapes.add_shape(MSO_SHAPE.CAN, x + Inches(2.65), y + Inches(1.52), Inches(1.12), Inches(1.58))
    set_fill(cyl, ACCENT_PALE)
    set_line(cyl, ACCENT, width=1.3)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(2.08), y + Inches(2.3), x + Inches(2.65), y + Inches(2.3))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2.2)
    add_text(slide, "volume", x + Inches(2.65), y + Inches(3.34), Inches(1.12), Inches(0.16), 8, MUTED, bold=True, align=PP_ALIGN.CENTER)


def draw_compose(slide):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Multi-services")
    items = [("web", 2.0, 1.35), ("api", 1.15, 2.6), ("db", 2.85, 2.6)]
    centers = []
    for label, px, py in items:
        bx, by = x + Inches(px), y + Inches(py)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(0.92), Inches(0.56))
        set_fill(box, ACCENT if label == "web" else WHITE)
        set_line(box, ACCENT, width=1.2)
        add_text(slide, label, bx, by + Inches(0.19), Inches(0.92), Inches(0.1), 8, WHITE if label == "web" else ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)
        centers.append((bx + Inches(0.46), by + Inches(0.28)))
    for a, b in [(0, 1), (0, 2), (1, 2)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, centers[a][0], centers[a][1], centers[b][0], centers[b][1])
        line.line.color.rgb = ACCENT_LIGHT
        line.line.width = Pt(2)
    file = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(1.45), y + Inches(3.78), Inches(1.84), Inches(0.32))
    set_fill(file, ACCENT_PALE)
    set_line(file, ACCENT_LIGHT, width=1)
    add_text(slide, "compose.yaml", x + Inches(1.45), y + Inches(3.88), Inches(1.84), Inches(0.08), 6.5, ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)


def draw_terminal(slide, command):
    x, y, _, _ = add_visual_panel(slide)
    add_visual_title(slide, x, y, "Commande")
    code = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.54), y + Inches(1.58), Inches(3.64), Inches(2.32))
    set_fill(code, ACCENT_DARK)
    no_line(code)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.54), y + Inches(1.58), Inches(3.64), Inches(0.38))
    set_fill(header, ACCENT)
    no_line(header)
    for i in range(3):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.78 + i * 0.2), y + Inches(1.72), Inches(0.08), Inches(0.08))
        set_fill(dot, WHITE)
        no_line(dot)
    add_text(slide, "$", x + Inches(0.84), y + Inches(2.42), Inches(0.22), Inches(0.22), 17, ACCENT_LIGHT, bold=True, font="Consolas")
    size = 18 if len(command) <= 30 else 15
    cmd = add_text(slide, command, x + Inches(1.12), y + Inches(2.42), Inches(2.74), Inches(0.56), size, WHITE, font="Consolas")
    cmd.text_frame.paragraphs[0].runs[0].font.name = "Consolas"


def draw_visual(slide, data, idx):
    if idx == 1:
        draw_title_visual(slide)
        return
    if data["command"]:
        draw_terminal(slide, data["command"])
        return
    visual = data.get("visual", "generic")
    if visual == "problem":
        draw_problem(slide)
    elif visual == "container":
        draw_container(slide)
    elif visual == "compare":
        draw_compare(slide)
    elif visual == "network":
        draw_network(slide)
    elif visual == "image":
        draw_image(slide)
    elif visual == "hub":
        draw_hub(slide)
    elif visual == "storage":
        draw_storage(slide)
    elif visual == "compose":
        draw_compose(slide)
    else:
        draw_generic(slide, section_label(data["section"]))


def build_deck():
    assert len(SLIDES) == 90, f"Nombre de slides source incorrect : {len(SLIDES)}"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        add_background(slide, data, idx)
        add_title(slide, data, idx)
        add_bullets(slide, data["bullets"], idx)
        draw_visual(slide, data, idx)

    assert len(prs.slides) == 90, f"Nombre de slides incorrect : {len(prs.slides)}"
    props = prs.core_properties
    props.title = "Docker - Concepts, utilisation et bonnes pratiques"
    props.subject = "Cours Docker haut de gamme avec design unique"
    props.author = "OpenAI Codex"
    props.keywords = "Docker, DevOps, cours, conteneurisation, Compose"
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build_deck()
    print(f"Présentation haut de gamme créée : {path.resolve()}")
    print(f"Nombre de slides : {len(SLIDES)}")
